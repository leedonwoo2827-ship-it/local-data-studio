"""
core/export.py — 보고서(마크다운)를 Excel / Word / PDF / PowerPoint 로 변환.

흐름: LLM 이 마크다운 보고서를 생성 → parse_blocks() 로 구조화 → 각 포맷 렌더러.
무거운 라이브러리(docx/pptx/reportlab)는 함수 내부에서 import 한다.
"""
from __future__ import annotations

import io
import os
import re

import pandas as pd


# ---------------------------------------------------------------------------
# 마크다운 파싱
# ---------------------------------------------------------------------------
def clean_inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)   # **굵게**
    text = re.sub(r"__(.+?)__", r"\1", text)
    text = text.replace("`", "").strip()
    return text


def parse_blocks(md: str) -> list[tuple[str, object]]:
    """마크다운 → [(type, content)]. type: h1/h2/h3/p/bullets/table."""
    blocks: list[tuple[str, object]] = []
    lines = md.splitlines()
    n = len(lines)
    i = 0
    para: list[str] = []
    bullets: list[str] = []

    def flush_para():
        nonlocal para
        if para:
            blocks.append(("p", " ".join(para).strip()))
            para = []

    def flush_bullets():
        nonlocal bullets
        if bullets:
            blocks.append(("bullets", bullets))
            bullets = []

    def cells(row: str) -> list[str]:
        return [c.strip() for c in row.strip().strip("|").split("|")]

    while i < n:
        s = lines[i].strip()
        if not s:
            flush_para(); flush_bullets(); i += 1; continue

        if s.startswith("|") and s.endswith("|"):
            flush_para(); flush_bullets()
            tbl = []
            while i < n and lines[i].strip().startswith("|"):
                tbl.append(lines[i].strip()); i += 1
            if tbl:
                headers = cells(tbl[0])
                rows = []
                for r in tbl[1:]:
                    compact = r.replace("|", "").replace(" ", "")
                    if compact and set(compact) <= set("-:"):
                        continue  # 구분선
                    rows.append(cells(r))
                blocks.append(("table", (headers, rows)))
            continue

        if s.startswith("### "):
            flush_para(); flush_bullets(); blocks.append(("h3", clean_inline(s[4:]))); i += 1; continue
        if s.startswith("## "):
            flush_para(); flush_bullets(); blocks.append(("h2", clean_inline(s[3:]))); i += 1; continue
        if s.startswith("# "):
            flush_para(); flush_bullets(); blocks.append(("h1", clean_inline(s[2:]))); i += 1; continue
        if s.startswith(("- ", "* ")):
            flush_para(); bullets.append(clean_inline(s[2:])); i += 1; continue
        m = re.match(r"^\d+\.\s+(.*)", s)
        if m:
            flush_para(); bullets.append(clean_inline(m.group(1))); i += 1; continue

        para.append(clean_inline(s)); i += 1

    flush_para(); flush_bullets()
    return blocks


def _group_sections(blocks: list[tuple[str, object]]):
    """h1/h2 기준으로 (제목, [blocks]) 섹션 묶음 — PPTX 슬라이드용."""
    sections: list[tuple[str, list]] = []
    title = "보고서"
    cur: list = []
    for t, c in blocks:
        if t in ("h1", "h2"):
            if cur or sections:
                sections.append((title, cur))
            title = c if isinstance(c, str) else "보고서"
            cur = []
        else:
            cur.append((t, c))
    sections.append((title, cur))
    return sections


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------
def to_excel(report_md: str, kpi_df: pd.DataFrame | None,
             data_df: pd.DataFrame | None, *, max_data_rows: int = 2000) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    from openpyxl.utils.dataframe import dataframe_to_rows

    wb = Workbook()

    # 1) 보고서 본문
    ws = wb.active
    ws.title = "보고서"
    ws.column_dimensions["A"].width = 110
    r = 1
    for t, c in parse_blocks(report_md):
        if t in ("h1", "h2", "h3"):
            cell = ws.cell(row=r, column=1, value=c)
            cell.font = Font(bold=True, size={"h1": 16, "h2": 13, "h3": 11}[t])
            r += 1
        elif t == "bullets":
            for it in c:
                ws.cell(row=r, column=1, value=f"• {it}").alignment = Alignment(wrap_text=True)
                r += 1
        elif t == "table":
            headers, rows = c
            for j, h in enumerate(headers, start=1):
                hc = ws.cell(row=r, column=j, value=h)
                hc.font = Font(bold=True)
                hc.fill = PatternFill("solid", fgColor="D9E1F2")
            r += 1
            for row in rows:
                for j, v in enumerate(row, start=1):
                    ws.cell(row=r, column=j, value=v)
                r += 1
        else:  # p
            ws.cell(row=r, column=1, value=c).alignment = Alignment(wrap_text=True)
            r += 1
        r += 0

    # 2) KPI 현황
    if kpi_df is not None and not kpi_df.empty:
        ws2 = wb.create_sheet("KPI현황")
        for row in dataframe_to_rows(kpi_df, index=False, header=True):
            ws2.append(row)
        for cell in ws2[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="D9E1F2")
        # 신호 색상
        if "신호" in list(kpi_df.columns):
            sig_idx = list(kpi_df.columns).index("신호") + 1
            for row in range(2, ws2.max_row + 1):
                v = ws2.cell(row=row, column=sig_idx).value
                if v == "Green":
                    ws2.cell(row=row, column=sig_idx).fill = PatternFill("solid", fgColor="C6EFCE")
                elif v == "Red":
                    ws2.cell(row=row, column=sig_idx).fill = PatternFill("solid", fgColor="FFC7CE")

    # 3) 원본 데이터
    if data_df is not None and not data_df.empty:
        ws3 = wb.create_sheet("데이터")
        for row in dataframe_to_rows(data_df.head(max_data_rows), index=False, header=True):
            ws3.append(row)
        for cell in ws3[1]:
            cell.font = Font(bold=True)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------
def to_docx(report_md: str) -> bytes:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    for t, c in parse_blocks(report_md):
        if t == "h1":
            doc.add_heading(c, level=0)
        elif t == "h2":
            doc.add_heading(c, level=1)
        elif t == "h3":
            doc.add_heading(c, level=2)
        elif t == "bullets":
            for it in c:
                doc.add_paragraph(it, style="List Bullet")
        elif t == "table":
            headers, rows = c
            table = doc.add_table(rows=1, cols=len(headers))
            try:
                table.style = "Light Grid Accent 1"
            except Exception:
                pass
            for j, h in enumerate(headers):
                cell = table.rows[0].cells[j]
                cell.text = h
                for p in cell.paragraphs:
                    for run in p.runs:
                        run.bold = True
            for row in rows:
                rc = table.add_row().cells
                for j, v in enumerate(row):
                    if j < len(rc):
                        rc[j].text = str(v)
        else:
            doc.add_paragraph(c)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF (.pdf) — 한글 폰트 임베딩
# ---------------------------------------------------------------------------
def _register_korean_font() -> str:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    candidates = [
        r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\malgunsl.ttf",
        r"C:\Windows\Fonts\gulim.ttc",
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
        "/Library/Fonts/AppleGothic.ttf",
    ]
    for p in candidates:
        if os.path.exists(p):
            try:
                pdfmetrics.registerFont(TTFont("KFont", p))
                return "KFont"
            except Exception:
                continue
    return "Helvetica"  # 한글 미표시 폴백


def to_pdf(report_md: str) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle, ListFlowable, ListItem)
    from reportlab.lib.styles import ParagraphStyle

    font = _register_korean_font()
    styles = {
        "h1": ParagraphStyle("h1", fontName=font, fontSize=18, leading=24, spaceAfter=10),
        "h2": ParagraphStyle("h2", fontName=font, fontSize=14, leading=20, spaceBefore=8, spaceAfter=6),
        "h3": ParagraphStyle("h3", fontName=font, fontSize=12, leading=16, spaceBefore=6, spaceAfter=4),
        "p":  ParagraphStyle("p",  fontName=font, fontSize=10, leading=15, spaceAfter=4),
    }

    buf = io.BytesIO()
    docp = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=18 * mm, rightMargin=18 * mm,
                            topMargin=18 * mm, bottomMargin=18 * mm)
    story = []
    for t, c in parse_blocks(report_md):
        if t in ("h1", "h2", "h3"):
            story.append(Paragraph(c, styles[t]))
        elif t == "bullets":
            items = [ListItem(Paragraph(it, styles["p"])) for it in c]
            story.append(ListFlowable(items, bulletType="bullet", start="•"))
            story.append(Spacer(1, 4))
        elif t == "table":
            headers, rows = c
            data = [headers] + rows
            tbl = Table(data, repeatRows=1)
            tbl.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, -1), font),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9E1F2")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F6FB")]),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 6))
        else:
            story.append(Paragraph(c, styles["p"]))

    docp.build(story)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------
def to_pptx(report_md: str, *, title: str = "보고서", subtitle: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Pt, Inches

    blocks = parse_blocks(report_md)
    prs = Presentation()

    # 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    blank = prs.slide_layouts[6]
    content_layout = prs.slide_layouts[1]

    def add_table_slide(headers, rows, sec_title):
        s = prs.slides.add_slide(blank)
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tx.text_frame.text = sec_title
        tx.text_frame.paragraphs[0].font.size = Pt(22)
        tx.text_frame.paragraphs[0].font.bold = True
        nrows, ncols = len(rows) + 1, max(1, len(headers))
        gt = s.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.1),
                                Inches(9), Inches(0.4 * nrows)).table
        for j, h in enumerate(headers):
            gt.cell(0, j).text = str(h)
        for i, row in enumerate(rows, start=1):
            for j in range(ncols):
                gt.cell(i, j).text = str(row[j]) if j < len(row) else ""

    for sec_title, body in _group_sections(blocks):
        # 불릿/문단 슬라이드
        bullet_items = []
        for t, c in body:
            if t == "bullets":
                bullet_items.extend(c)
            elif t == "p":
                bullet_items.append(c)
        if bullet_items:
            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = sec_title
            tf = s.placeholders[1].text_frame
            tf.clear()
            for k, it in enumerate(bullet_items):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.text = it
                p.font.size = Pt(16)
        # 표는 별도 슬라이드
        for t, c in body:
            if t == "table":
                headers, rows = c
                add_table_slide(headers, rows, f"{sec_title} — 표")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
